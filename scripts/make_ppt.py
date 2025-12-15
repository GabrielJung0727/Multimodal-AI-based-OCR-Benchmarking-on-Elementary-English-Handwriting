"""
Generate a Korean PPTX deck from the mock results.
- Reads graphs from outputs/graphs/*.png
- Saves to multimodal_ocr_benchmark.pptx
Requires: pip install python-pptx
"""

from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
GRAPH_DIR = ROOT / "outputs" / "graphs"


def add_title(slide, text: str) -> None:
    title = slide.shapes.title
    title.text = text
    for paragraph in title.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.font.name = "Arial"
        paragraph.font.size = Pt(34)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(30, 58, 138)  # #1e3a8a


def add_content(
    slide,
    text_lines,
    *,
    left: float = 1.0,
    top: float = 1.5,
    width: float = 8.5,
    height: float = 5.0,
) -> None:
    tx_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(text_lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = line
        p.font.name = "Arial"
        p.font.size = Pt(19)
        p.space_after = Pt(10)
        p.level = 0


def add_image(slide, img_name: str, *, left: float, top: float, height: float) -> None:
    img_path = GRAPH_DIR / img_name
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(left), Inches(top), height=Inches(height))
    else:
        tx_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(4), Inches(1))
        tx_box.text = f"[이미지 없음: {img_path}]"


def create_presentation() -> None:
    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content

    # 1. 프로젝트 개요
    slide = prs.slides.add_slide(layout)
    add_title(slide, "프로젝트 개요")
    add_content(
        slide,
        [
            "목표: 필기체 영어 이미지에서 6개 멀티모달 OCR/비전-언어 모델 벤치마킹.",
            "모델: Gemini, Qwen, LLaMA, Clova, GPT, Claude.",
            "지표: OCR 정확도, 오류율, 신뢰도 안정성.",
        ],
    )

    # 2. 데이터 소스 및 방법론
    slide = prs.slides.add_slide(layout)
    add_title(slide, "데이터 소스 및 방법론")
    add_content(
        slide,
        [
            "데이터셋:",
            "- 소스: data/kakao (필기체/사진 영어 텍스트)",
            "- 모든 이미지를 단일 매니페스트로 구성, 전 모델 동일 적용",
            "",
            "평가 방식:",
            "- 외부 GT 라벨 없이 모델 출력과 신뢰도로 비교",
            "- 결과/메트릭은 outputs/ai_judgements/, outputs/metrics/에 저장",
        ],
    )

    # 3. 학습 설정
    slide = prs.slides.add_slide(layout)
    add_title(slide, "학습 설정")
    add_content(
        slide,
        [
            "시뮬레이션: 모델별 30 에폭 (파인튜닝 스타일).",
            "순위 유지: Gemini > Qwen ≈ LLaMA > Clova > GPT ≈ Claude.",
            "곡선 변동성: 에폭별 노이즈/지터를 넣어 비균질하게 생성.",
            "산출물: outputs/ai_judgements/*.jsonl, outputs/metrics/*.csv",
        ],
    )

    # 4. 에폭별 정확도
    slide = prs.slides.add_slide(layout)
    add_title(slide, "에폭별 정확도")
    add_content(
        slide,
        [
            "Gemini가 선두, Qwen/LLaMA가 근소 추격.",
            "Clova/GPT는 완만 상승, Claude는 소폭 뒤처짐.",
        ],
        width=4.5,
    )
    add_image(slide, "accuracy_per_epoch.png", left=5.0, top=2.0, height=4.5)

    # 5. 에폭별 오류율
    slide = prs.slides.add_slide(layout)
    add_title(slide, "에폭별 오류율")
    add_content(
        slide,
        [
            "오류 감소: Gemini/Qwen 가장 빠름, LLaMA 유사 추세.",
            "Clova/GPT는 높은 구간에서 완만 정체, Claude가 잔여 오류 가장 높음.",
        ],
        width=4.5,
    )
    add_image(slide, "error_rate_per_epoch.png", left=5.0, top=2.0, height=4.5)

    # 6. Gemini 대비 정확도 격차
    slide = prs.slides.add_slide(layout)
    add_title(slide, "Gemini 대비 정확도 격차")
    add_content(
        slide,
        [
            "각 에폭에서 Gemini와의 포인트 격차 시각화.",
            "Qwen/LLaMA는 좁은 격차, Clova/GPT/Claude는 더 큰 격차.",
        ],
        width=4.5,
    )
    add_image(slide, "accuracy_gap.png", left=5.0, top=2.0, height=4.5)

    # 7. 최종 정확도
    slide = prs.slides.add_slide(layout)
    add_title(slide, "최종 정확도 (30 에폭 후)")
    add_content(
        slide,
        [
            "대략적 최종 정확도: Gemini ~96%, Qwen/LLaMA ~90%, Clova ~86%, GPT ~82%, Claude ~81%.",
            "명확한 계층 구조를 유지하되, 모델 간 격차는 작게 분산.",
        ],
        width=4.5,
    )
    add_image(slide, "final_accuracy.png", left=5.0, top=2.0, height=4.5)

    # 8. 신뢰도 분포
    slide = prs.slides.add_slide(layout)
    add_title(slide, "신뢰도(Confidence) 분포")
    add_content(
        slide,
        [
            "신뢰도 중앙값: Gemini/Qwen > LLaMA > Clova > GPT/Claude.",
            "박스 폭이 좁을수록 일관성이 높음.",
        ],
        width=4.5,
    )
    add_image(slide, "confidence_box.png", left=5.0, top=2.0, height=4.5)

    out_file = ROOT / "multimodal_ocr_benchmark.pptx"
    prs.save(out_file)
    print(f"'{out_file.name}' 파일이 생성되었습니다.")


if __name__ == "__main__":
    create_presentation()
